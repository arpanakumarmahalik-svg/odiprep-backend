import fitz  # PyMuPDF
import os
from docx import Document
from pptx import Presentation

def extract_text_from_pdf(file_path: str) -> str:
    """Read all text from a PDF file"""
    text = ""
    doc = fitz.open(file_path)
    for page in doc:
        text += page.get_text()
    doc.close()
    return text.strip()

def extract_text_from_txt(file_path: str) -> str:
    """Read all text from a TXT file"""
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read().strip()

def extract_text_from_docx(file_path: str) -> str:
    """Read all text from a DOCX file"""
    doc = Document(file_path)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text.strip()

def extract_text_from_pptx(file_path: str) -> str:
    """Read all text from a PPT file"""
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

def extract_text(file_path: str) -> str:
    """
    Main function — detects file type and calls
    the right extractor automatically
    """
    try:
        extension = file_path.split(".")[-1].lower()

        if extension == "pdf":
            return extract_text_from_pdf(file_path)
        elif extension == "txt":
            return extract_text_from_txt(file_path)
        elif extension == "docx":
            return extract_text_from_docx(file_path)
        elif extension == "pptx":
            return extract_text_from_pptx(file_path)
        else:
            return ""
    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}")
        return ""